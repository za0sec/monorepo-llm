"""
Genera la presentacion del TP1 (Predecir BTR con Transformers) en .pptx,
editable despues en Google Slides.

Contenido tomado exclusivamente de ejercicio1/Notas.md y ejercicio2/
Notas.md + Experimentos.md -- no se agrega ninguna afirmacion, tecnica
o numero que no este en esas notas o en las clases citadas ahi.

Sigue el estilo de SIA-TP4/SIA-PCA/generate_pptx.py: pocas slides por
paso, texto + imagen, sin plantillas de diseño pesadas (para que se
pueda seguir editando comodo en Slides).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = os.path.dirname(os.path.abspath(__file__))
EJ1 = os.path.join(BASE, "ejercicio1")
EJ2 = os.path.join(BASE, "ejercicio2", "output")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0x4A, 0x6F, 0xC5)
DARK = RGBColor(0x33, 0x33, 0x33)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
GRAY = RGBColor(0x70, 0x70, 0x70)
BLANK = 6


def slide():
    return prs.slides.add_slide(prs.slide_layouts[BLANK])


def add_title(s, text, top=0.3, size=30):
    tx = s.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.3), Inches(0.9))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = DARK


def add_text(s, text, left, top, width, height, size=14, bold=False, color=DARK,
             align=None, italic=False):
    tx = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = color
        if align is not None:
            p.alignment = align
    return tx


def add_bullets(s, items, left, top, width, height, size=16, color=DARK):
    tx = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return tx


def add_table(s, headers, rows, left, top, width, height, size=12, col_widths=None):
    n_rows, n_cols = len(rows) + 1, len(headers)
    shape = s.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top),
                                Inches(width), Inches(height))
    table = shape.table
    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(width * w / total)
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = str(h)
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(size)
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(size)
                p.font.color.rgb = DARK
    return table


def add_params_box(s, params, left, top, width, height):
    """Caja gris con los parametros de la corrida, para slides con grafico."""
    tx = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Parámetros de la corrida:"
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    for item in params:
        pp = tf.add_paragraph()
        pp.text = f"– {item}"
        pp.font.size = Pt(11)
        pp.font.color.rgb = GRAY
    return tx


def add_conclusion_box(s, text, left, top, width, height):
    tx = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Conclusión: {text}"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = GREEN
    return tx


def add_image(s, path, left, top, height=None, width=None):
    kwargs = {}
    if height:
        kwargs["height"] = Inches(height)
    if width:
        kwargs["width"] = Inches(width)
    s.shapes.add_picture(path, Inches(left), Inches(top), **kwargs)


def section_slide(title, subtitle=""):
    s = slide()
    add_text(s, title, 0.5, 3.0, 12.3, 1.2, size=40, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(s, subtitle, 0.5, 4.2, 12.3, 0.6, size=18, color=DARK, align=PP_ALIGN.CENTER)


# ============================================================
# 1. TITULO
# ============================================================
s = slide()
add_text(s, "Predicción de Buy Through Rate (BTR)\ncon Transformers", 0.5, 2.3, 12.3, 1.8,
         size=40, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_text(s, "TP1 — Transformers", 0.5, 4.1, 12.3, 0.6, size=20, align=PP_ALIGN.CENTER)
add_text(s, "73.69 Large Language Models — 2026", 0.5, 6.8, 12.3, 0.4, size=14, color=BLUE,
         align=PP_ALIGN.CENTER)

# ============================================================
# 2. OBJETIVO
# ============================================================
s = slide()
add_title(s, "Objetivo")
add_bullets(s, [
    "Predecir el Buy Through Rate (BTR) en un e-commerce de supermercado: "
    "productos comprados / productos impresos en resultados de búsqueda.",
    "Dataset: supermarket_products.csv — eventos de búsqueda (compras, vistas, "
    "interacciones con filtros).",
    "El desarrollo debe contar con al menos un modelo basado en la arquitectura Transformer.",
    "Foco del TP: comprensión de la arquitectura y su aplicación a un problema concreto — "
    "se evalúa la justificación de decisiones y la comparación de alternativas, no solo "
    "el mejor resultado numérico.",
], 0.7, 1.5, 11.9, 4.5, size=18)
add_bullets(s, [
    "Ejercicio 1 — Formulación del problema y EDA",
    "Ejercicio 2 — Desarrollo del sistema",
    "Ejercicio 3 — Personalización (teórico)",
], 0.7, 5.6, 11.9, 1.6, size=16, color=BLUE)

section_slide("Ejercicio 1", "Formulación del problema y EDA")

# ============================================================
# Punto de partida
# ============================================================
s = slide()
add_title(s, "Punto de partida del dataset")
add_bullets(s, [
    "10.000 filas, 22 columnas.",
    "Cada fila NO es una búsqueda completa — es un producto puntual mostrado como "
    "resultado dentro de una búsqueda (query_id).",
    "2.012 búsquedas (query_id) en total, entre 1 y 8 productos impresos por búsqueda "
    "(~5 en promedio).",
    "Columnas cart y bought indican, por fila, si ESE producto en ESA búsqueda se "
    "agregó al carrito / se compró.",
], 0.7, 1.6, 11.9, 4.5, size=18)

# ============================================================
# Variable objetivo
# ============================================================
s = slide()
add_title(s, "Variable objetivo: bought a nivel fila")
add_bullets(s, [
    "El BTR de la consigna es una métrica agregada a nivel búsqueda, pero el dataset "
    "viene desagregado a nivel fila-producto-búsqueda.",
    "Variable objetivo elegida: bought, a nivel fila — probabilidad de que ESE producto "
    "puntual sea comprado.",
    "El BTR de una búsqueda sale de promediar bought (o las probabilidades predichas) "
    "agrupando por query_id.",
], 0.7, 1.5, 11.9, 3.0, size=18)
add_conclusion_box(s,
    "consistente con evaluar con PR-AUC/ROC-AUC (clasificación binaria) y con que no haga "
    "falta threshold — importa qué tan bien el modelo ordena los productos por probabilidad "
    "de compra, no un corte fijo.",
    0.7, 5.2, 11.9, 1.6)

# ============================================================
# cart vs bought
# ============================================================
s = slide()
add_title(s, "Relación cart y bought")
add_table(s, ["cart", "bought", "% filas"],
          [["True", "True", "13,01%"], ["True", "False", "17,06%"],
           ["False", "True", "0,00%"], ["False", "False", "69,93%"]],
          0.7, 1.5, 5.5, 2.5, size=14)
add_bullets(s, [
    "bought=True implica siempre cart=True (0% de compras sin carrito).",
    "No al revés: de los que llegaron al carrito (30,07%), solo el 43,3% terminó "
    "comprado.",
], 6.6, 1.6, 6.0, 2.2, size=16)
add_conclusion_box(s,
    "cart no puede usarse como feature — al momento de predecir (búsqueda mostrada) "
    "todavía no se sabe si se va a agregar al carrito. Sería data leakage.",
    0.7, 5.2, 11.9, 1.4)

# ============================================================
# Balance de clases
# ============================================================
s = slide()
add_title(s, "Balance de clases")
add_bullets(s, [
    "bought = True: 13,01% de las filas.",
    "cart = True: 30,07% de las filas.",
    "Desbalance moderado (no extremo).",
], 0.7, 1.6, 11.9, 2.2, size=20)
add_conclusion_box(s,
    "PR-AUC más informativo que ROC-AUC acá: el ROC-AUC puede verse artificialmente "
    "alto por la cantidad de negativos (no-compras) fáciles de acertar.",
    0.7, 4.5, 11.9, 1.4)

# ============================================================
# Nulos en allergens
# ============================================================
s = slide()
add_title(s, "Nulos en allergens (única columna con nulos, 44,55%)")
add_bullets(s, [
    "Nulo siempre (100%): Household, Meat, Personal Care, Produce — no aplica declarar "
    "alérgenos alimentarios.",
    "Nunca nulo (0%): Bakery, Dairy, Seafood.",
    "Parcial: Beverages (50,8%), Baby (34,0%), Frozen (25,2%), Pantry (25,1%), Snacks (18,6%).",
    "No es aleatorio: depende 100% de category.",
], 0.7, 1.5, 11.9, 3.4, size=17)
add_conclusion_box(s,
    "el nulo codifica información real (\"sin alérgenos declarados\"), no un dato perdido "
    "— tratarlo como categoría propia al encodear, no imputar ni descartar.",
    0.7, 5.1, 11.9, 1.4)

s = slide()
add_title(s, "allergens no es consistente con ingredients")
add_bullets(s, [
    "Se probó si el alérgeno declarado aparece como substring en la lista de ingredientes "
    "de la misma fila.",
    "La misma combinación exacta de ingredientes aparece unas veces con un alérgeno y "
    "otras con otro (ej. \"Prepared ingredients, Spices, Salt\" con Wheat y con Soy).",
    "% de coincidencia: 0% para Fish, Peanuts, Shellfish, Soy, Tree nuts; parcial (55-68%) "
    "solo para Wheat/Milk (por casualidad de palabras como \"Wheat flour\").",
], 0.7, 1.6, 11.9, 3.6, size=17)
add_conclusion_box(s,
    "allergens se generó de forma independiente de ingredients (probablemente condicionado "
    "solo por category, igual que los nulos) — no se puede usar una para derivar la otra.",
    0.7, 5.3, 11.9, 1.4)

# ============================================================
# Panorama de columnas
# ============================================================
s = slide()
add_title(s, "Panorama general de columnas (22 en total)")
add_bullets(s, [
    "Texto libre / casi-únicos: title (9910), description (9112), dimensions_in (9864).",
    "Categóricas de pocas opciones: category/filter_category (12), storage_type/"
    "filter_storage_type (3), unit_of_measure (5), brand (15), country_of_origin (10), "
    "allergens (7).",
    "Numéricas continuas: price, filter_price_min, filter_price_max, net_weight_oz, "
    "nutrition_score.",
    "Identificadores / contexto de búsqueda: query_id, timestamp.",
    "Booleanas: cart (leakage), bought (target).",
    "filter_category y filter_storage_type: iguales a category/storage_type en el 100% "
    "de las filas — descartadas por redundancia.",
], 0.7, 1.5, 11.9, 5.3, size=15)

# ============================================================
# dimensions_in
# ============================================================
s = slide()
add_title(s, "dimensions_in: sin estándar de envase")
add_bullets(s, [
    "Formato: \"3.3 x 4.0 x 4.1\\\"\" (largo x ancho x alto, pulgadas).",
    "Hipótesis: no hay estándar de envase, cada producto tiene su propia medida.",
    "Volumen vs net_weight_oz: correlación 0,38 (moderada-baja, no exacta).",
    "Solo 136 de 10.000 filas con dimensions_in exactamente repetido.",
], 0.7, 1.6, 11.9, 3.4, size=18)
add_conclusion_box(s,
    "no tratar como categórica (quedaría con ~9864 categorías) — parsear a 3 numéricas "
    "(largo, ancho, alto) o resumir en volumen.",
    0.7, 5.3, 11.9, 1.4)

# ============================================================
# package_size / unit_of_measure
# ============================================================
s = slide()
add_title(s, "package_size / unit_of_measure vs net_weight_oz")
add_bullets(s, [
    "unit_of_measure: oz, ct, lb, gal, fl oz. package_size: texto \"10 oz\", \"12 ct\", etc.",
    "Aún con unit_of_measure='oz', el número de package_size casi nunca coincide exacto "
    "con net_weight_oz (0,4% de match exacto).",
    "Para ct/lb/gal/fl oz no hay fórmula de conversión limpia (mismo package_size, "
    "distinto net_weight_oz).",
], 0.7, 1.6, 11.9, 3.2, size=17)
add_conclusion_box(s,
    "package_size y unit_of_measure son redundantes con net_weight_oz (ya normalizado a "
    "onzas) — se descartan como features.",
    0.7, 5.1, 11.9, 1.4)

# ============================================================
# ingredients
# ============================================================
s = slide()
add_title(s, "ingredients: 190 valores → 12 combinaciones reales")
add_bullets(s, [
    "nunique crudo: 190. Ordenando los ingredientes dentro de cada fila (ignorando "
    "orden): solo 12 combinaciones reales.",
    "Vocabulario total: 21 ingredientes individuales distintos.",
], 0.7, 1.6, 11.9, 2.4, size=18)
add_conclusion_box(s,
    "no tratar como texto libre ni 190 categorías. Dos caminos: multi-hot sobre los 21 "
    "ingredientes, o categórica de 12 combinaciones canónicas.",
    0.7, 4.3, 11.9, 1.4)

# ============================================================
# price vs filtros - LA CAMPANA
# ============================================================
s = slide()
add_title(s, "price vs. filter_price_min / filter_price_max")
add_image(s, os.path.join(EJ1, "price_position_bought.png"), 0.5, 1.4, height=4.6)
add_params_box(s, [
    "pos = (price − filter_price_min) / (filter_price_max − filter_price_min)",
    "bins de 0.2 en 0.2, sobre 10.000 filas",
], 7.3, 1.5, 5.5, 1.6)
add_bullets(s, [
    "price siempre cae dentro de [filter_price_min, filter_price_max] — 100% de las filas.",
    "0,0-0,2: 7,49% (n=1935)  ·  0,4-0,6: 17,56% (n=2569)  ·  0,8-1,0: 9,01% (n=677)",
], 7.3, 3.3, 5.5, 2.2, size=13)
add_conclusion_box(s,
    "forma de campana: se compra menos lo más barato y lo más caro del rango buscado, más "
    "lo intermedio. Feature con señal real — usar la posición relativa, no solo price.",
    0.5, 6.2, 12.3, 1.1)

# ============================================================
# Encoding categoricas nominales
# ============================================================
s = slide()
add_title(s, "Encoding de categóricas nominales")
add_bullets(s, [
    "category (12), storage_type (3), brand (15), country_of_origin (10), allergens (7) "
    "— todas nominales, sin orden natural.",
    "One-hot encoding (sugerido por la consigna): N columnas binarias, evita imponer "
    "orden/magnitud falsos.",
    "country_of_origin muy desbalanceada: 75% (7500/10000) es \"United States\", el resto "
    "en 9 países con 245-331 filas c/u.",
], 0.7, 1.5, 11.9, 3.4, size=17)
add_conclusion_box(s,
    "riesgo del desbalance en one-hot: categorías con pocos ejemplos dan poca evidencia "
    "para un peso confiable — puede ajustarse a ruido en vez de señal real.",
    0.7, 5.1, 11.9, 1.5)

# ============================================================
# country_of_origin sin señal
# ============================================================
s = slide()
add_title(s, "country_of_origin vs. bought")
add_table(s, ["país", "% bought", "n"],
          [["Vietnam", "16,61%", "271"], ["Italy", "14,55%", "275"], ["Peru", "13,88%", "245"],
           ["Mexico", "13,77%", "305"], ["United States", "12,99%", "7500"],
           ["Canada", "12,73%", "267"], ["New Zealand", "12,32%", "284"],
           ["Spain", "12,08%", "331"], ["Chile", "11,11%", "261"], ["Thailand", "10,73%", "261"]],
          0.7, 1.4, 6.5, 4.8, size=12)
add_text(s, "Promedio general: 13,01%", 7.6, 1.6, 5.2, 0.5, size=14, italic=True, color=GRAY)
add_bullets(s, [
    "Rango angosto (10,7%-16,6%), sin país que se dispare.",
    "9 países no-US con 245-331 filas c/u: diferencias de 2-3 puntos son compatibles con "
    "ruido (margen esperado por azar ±4 puntos con n≈245).",
], 7.6, 2.2, 5.2, 3.0, size=15)
add_conclusion_box(s, "sin señal clara — candidata a simplificar o evaluar en ablación.",
                    7.6, 5.5, 5.2, 1.4)

# ============================================================
# nutrition_score sin señal
# ============================================================
s = slide()
add_title(s, "nutrition_score vs. bought")
add_bullets(s, [
    "Correlación con bought prácticamente nula: -0,019.",
    "Tasa de bought por bins de 20 puntos: entre 12,1% y 14,1% (promedio general 13,01%), "
    "sin tendencia — ni \"más nutritivo vende más\" ni lo contrario.",
], 0.7, 1.7, 11.9, 2.4, size=19)
add_conclusion_box(s,
    "sin señal en el cruce univariado — igual que country_of_origin, candidata a "
    "incluir/excluir en el estudio de ablación (podría aportar en combinación con otras "
    "features).",
    0.7, 4.3, 11.9, 1.6)

# ============================================================
# timestamp y query_id sin señal
# ============================================================
s = slide()
add_title(s, "timestamp y query_id: sin señal encontrada")
add_bullets(s, [
    "timestamp: se probó si la hora del día predice bought (hipótesis: alcohol/Beverages "
    "a la noche, desayuno/Bakery a la mañana). Tasa por hora salta erráticamente entre "
    "horas consecutivas (ej. Beverages 21h=21,3% pero 20h=10,5%) — ruido, con solo "
    "30-58 filas por hora dentro de cada categoría.",
    "query_id: el ID en sí no debe usarse (no generaliza). La feature derivada "
    "\"n_resultados de la búsqueda\" (2 a 8) tampoco mostró señal (tasa entre 12-15% "
    "sin importar el tamaño).",
], 0.7, 1.6, 11.9, 4.0, size=17)
add_conclusion_box(s,
    "ambas descartadas como features. query_id sigue siendo necesario como clave para "
    "agrupar filas y calcular el BTR agregado por búsqueda.",
    0.7, 5.9, 11.9, 1.2)

# ============================================================
# Normalizacion numerica
# ============================================================
s = slide()
add_title(s, "Preprocesamiento numéricas: normalización z-score")
add_image(s, os.path.join(EJ1, "normalizacion_boxplot.png"), 0.5, 1.4, height=4.5)
add_params_box(s, [
    "price, net_weight_oz, nutrition_score",
    "z-score: (x − media) / desvío",
], 7.3, 1.5, 5.5, 1.3)
add_bullets(s, [
    "Escalas muy distintas: price (1,2-35), net_weight_oz (2,8-155), nutrition_score (0-99).",
    "Min-max descartado: muy sensible a outliers. net_weight_oz tiene cola larga "
    "(media 28,7, desvío 29) — con min-max la mayoría quedaría apretada cerca de 0.",
    "Consistencia con TP4 (SIA): mismo criterio de z-score usado en PCA/Kohonen/Oja.",
], 7.3, 3.0, 5.5, 2.8, size=13)

# ============================================================
# HALLAZGO ESTRELLA: tag de reputacion
# ============================================================
s = slide()
add_title(s, "Hallazgo: tag de reputación oculto en title", size=28)
add_image(s, os.path.join(EJ1, "reputation_tag_bought.png"), 0.4, 1.3, height=5.6)
add_params_box(s, [
    "regex: \\(([^)]+)\\)\\s*$ sobre title",
    "cruce con bought, 10.000 filas",
], 8.0, 1.5, 4.8, 1.3)
add_bullets(s, [
    "Customer Favorite: 67,75% (n=493)",
    "Best Seller: 65,74% (n=470)",
    "Top Rated: 62,71% (n=472)",
    "#1 Pick: 62,50% (n=496)",
    "Well Reviewed / Shopper Favorite /",
    "Highly Rated / Popular Choice: 1,9-3,8%",
    "Las otras 11 (incl. sin tag): 0,00%",
], 8.0, 3.0, 4.8, 3.5, size=13)
add_conclusion_box(s,
    "señal más fuerte de todo el EDA (base 13,01%). Se decide NO parsearla — dejar title "
    "crudo para que el Transformer la aprenda vía atención (foco: comprensión de la "
    "arquitectura, no un atajo).",
    0.4, 6.6, 12.5, 0.8)

# ============================================================
# Tokenizacion
# ============================================================
s = slide()
add_title(s, "Tokenización de title/description")
add_image(s, os.path.join(EJ1, "vocab_freq_rank.png"), 0.5, 1.4, height=4.4)
add_params_box(s, [
    "title + description concatenados, 10.000 filas",
    "minúsculas, split por no-alfanumérico, números como token propio",
], 7.3, 1.5, 5.5, 1.3)
add_bullets(s, [
    "357.419 tokens totales (35,74 promedio por fila).",
    "Vocabulario único: 410 palabras. Frecuencia mínima: 24 (0 hapax).",
    "Texto generado por plantilla (frases fijas con sustituciones acotadas), vocabulario "
    "cerrado, no abierto.",
], 7.3, 3.0, 5.5, 2.6, size=13)
add_conclusion_box(s,
    "tokenización por palabra, no BPE — BPE resuelve vocabulario abierto/OOV, problemas "
    "que acá no existen. <UNK> y <PAD> reservados igual, por las dudas.",
    0.5, 6.1, 12.3, 1.2)

# ============================================================
# Lista final de features
# ============================================================
s = slide()
add_title(s, "Lista final de features")
add_bullets(s, [
    "Descartadas (leakage): cart.",
    "Descartadas (redundantes): filter_category, filter_storage_type, package_size, "
    "unit_of_measure.",
    "Numéricas (z-score): net_weight_oz, posición relativa de price, dimensiones "
    "parseadas.",
    "One-hot: category, storage_type, brand, allergens.",
    "Multi-valor: ingredients (multi-hot o 12 combinaciones).",
    "Dudosas (a definir en ablación): country_of_origin, nutrition_score.",
    "Descartadas sin señal: timestamp, query_id (como feature).",
    "Para el Transformer: title, description — tokenización por palabra.",
], 0.7, 1.5, 11.9, 5.3, size=16)

section_slide("Ejercicio 2", "Desarrollo del sistema")

# ============================================================
# Split - por que agrupar
# ============================================================
s = slide()
add_title(s, "Split train / valid / test: por qué agrupar por query_id")
add_bullets(s, [
    "Las filas están agrupadas por búsqueda — comparten contexto (misma categoría, "
    "mismo rango de precio filtrado).",
    "Si se particiona fila por fila al azar, el modelo vería ese contexto exacto en "
    "train y test — no refleja el caso real (predecir búsquedas nuevas).",
    "Partición por query_id completo: todas las filas de una búsqueda van al mismo split.",
], 0.7, 1.6, 11.9, 3.2, size=18)
add_bullets(s, [
    "train: ajusta los pesos.",
    "valid: compara configuraciones durante la iteración.",
    "test: se toca una sola vez, al final, para el número reportado.",
], 0.7, 5.0, 11.9, 2.0, size=16, color=BLUE)

# ============================================================
# Split - proporciones y estratificacion
# ============================================================
s = slide()
add_title(s, "Split: 70/15/15 estratificado por tasa de bought")
add_image(s, os.path.join(EJ2, "split_balance.png"), 0.4, 1.4, height=4.3)
add_params_box(s, [
    "train_test_split (sklearn) en 2 pasos, stratify=franja de bought por query",
    "franjas: 0% / 1-33% / 34-100%",
    "random_state=42",
], 7.5, 1.5, 5.3, 1.7)
add_table(s, ["split", "n queries", "n filas", "tasa bought"],
          [["train", "1408", "7012", "12,94%"], ["valid", "302", "1498", "12,88%"],
           ["test", "302", "1490", "13,49%"]],
          7.5, 3.4, 5.3, 2.0, size=13)
add_conclusion_box(s,
    "trade-off: train grande ayuda con el desbalance (13%), pero valid/test chicos dan "
    "métricas ruidosas (ya visto con country_of_origin). Se prioriza confiabilidad: "
    "70/15/15, tasa global 13,01% respetada en las 3 particiones.",
    0.4, 5.9, 12.5, 1.4)

# ============================================================
# Arquitectura: Encoder-only
# ============================================================
s = slide()
add_title(s, "Arquitectura del bloque de texto: Encoder-only")
add_bullets(s, [
    "Encoder: Multi-Head Self-Attention (+ residual + Layer Norm) → MLP (+ residual + "
    "Layer Norm). Decoder agrega Cross-Attention y máscara (no ve tokens futuros).",
    "Encoder-Decoder completo: para secuencia-a-secuencia (ej. traducción) — no aplica.",
    "Decoder-only (como GPT, visto en la demo): generación autoregresiva — no es nuestro "
    "caso, no queremos generar texto.",
    "Encoder-only (como BERT): da un embedding/representación, sin máscara — el título "
    "ya está completo, no hay nada \"futuro\" que ocultar.",
], 0.7, 1.5, 11.9, 4.6, size=16)
add_conclusion_box(s,
    "se elige Encoder-only: la tarea es clasificar (bought), no generar texto.",
    0.7, 6.3, 11.9, 0.9)

# ============================================================
# Arquitectura elegida + fusion tardia
# ============================================================
s = slide()
add_title(s, "Arquitectura elegida: fusión tardía")
add_bullets(s, [
    "Bloque de texto: Encoder-only sobre title+description, con positional encoding "
    "senoidal (el visto en clase) y mean-pooling sobre los tokens no-pad para resumir "
    "la secuencia en un vector.",
    "Bloque tabular: vector de 75 columnas (numéricas z-score + one-hot + multi-hot).",
    "Fusión tardía: se concatena el vector de texto con el vector tabular, sigue por "
    "una capa densa hasta 1 logit de salida.",
], 0.7, 1.5, 11.9, 3.6, size=17)
add_params_box(s, [
    "Configuración base: d_model=64, 4 heads, 2 encoders apilados,",
    "MLP interno=128, dropout=0.1",
], 0.7, 5.3, 11.9, 1.2)
add_conclusion_box(s,
    "\"sacar el Transformer\" es desconectar una rama entera — ablación limpia. Diales "
    "citados en clase: cantidad de heads, encoders apilados, dimensión del MLP, d_model.",
    0.7, 6.5, 11.9, 0.8)

# ============================================================
# Setup de entrenamiento
# ============================================================
s = slide()
add_title(s, "Setup de entrenamiento")
add_bullets(s, [
    "Optimizador: Adam, lr=1e-3.",
    "Loss: BCEWithLogitsLoss.",
    "Batch size: 128. Épocas: 20.",
    "Sin ponderación de clases ni técnica específica para el desbalance (13% positivos) "
    "— no se vio nada puntual en clase para esto; se confía en PR-AUC/ROC-AUC.",
    "3 semillas por configuración (media ± desvío) — recomendación de consigna.VTT de "
    "promediar varias corridas.",
], 0.7, 1.6, 11.9, 4.6, size=17)

# ============================================================
# Experimento 1 vs 2
# ============================================================
s = slide()
add_title(s, "¿Aporta el Transformer? Baseline vs. fusión")
add_image(s, os.path.join(EJ2, "experiment_comparison.png"), 0.4, 1.4, height=4.3)
add_table(s, ["experimento", "PR-AUC (valid)", "ROC-AUC (valid)"],
          [["1. Baseline tabular", "0,178 ± 0,000", "0,580 ± 0,002"],
           ["2. Fusión (Encoder-only + tabular)", "0,732 ± 0,017", "0,965 ± 0,003"]],
          7.3, 1.6, 5.5, 1.6, size=12)
add_params_box(s, [
    "Baseline: TabularMLPBaseline, capas densas sobre 75 cols",
    "Fusión: EncoderFusionModel, config base (d_model=64)",
    "3 semillas c/u",
], 7.3, 3.4, 5.5, 1.6)
add_conclusion_box(s,
    "salto de +0,55 PR-AUC / +0,39 ROC-AUC. Confirma con números el hallazgo del EDA: la "
    "señal del tag de reputación domina la predicción.",
    0.4, 6.1, 12.5, 1.2)

# ============================================================
# Curvas de entrenamiento
# ============================================================
s = slide()
add_title(s, "Curvas de entrenamiento")
add_image(s, os.path.join(EJ2, "training_curves.png"), 1.5, 1.4, height=5.0)
add_conclusion_box(s,
    "baseline tabular mejora lento y monótono, sin overfitting (poca señal para "
    "explotar). Fusión converge muy rápido (la mayor parte de la mejora en la primera "
    "época) — consistente con un patrón simple y regular de detectar (el tag).",
    0.7, 6.5, 11.9, 0.9)

# ============================================================
# Overfitting / underfitting
# ============================================================
s = slide()
add_title(s, "Overfitting y underfitting (train vs. valid)")
add_image(s, os.path.join(EJ2, "overfitting_diagnosis.png"), 0.3, 1.3, height=4.2)
add_table(s, ["corrida", "PR-AUC train", "PR-AUC valid", "brecha"],
          [["tabular (media)", "~0,277", "~0,176", "+0,10"],
           ["fusión (media)", "~0,921", "~0,728", "+0,19"]],
          7.6, 1.5, 5.2, 1.4, size=12)
add_params_box(s, [
    "Métricas medidas también sobre train en modo eval",
    "(sin dropout, pesos fijos al final de la época)",
], 7.6, 3.1, 5.2, 1.1)
add_conclusion_box(s,
    "tabular: underfitting claro (las dos bajas y pegadas — no hay casi nada que "
    "aprender). Fusión: sí hay overfitting (brecha ~0,19), pero \"benigno\" — valid no "
    "degrada, sigue subiendo/amesetada hasta la época 20.",
    0.3, 5.0, 12.5, 2.0)

# ============================================================
# Interpretabilidad
# ============================================================
s = slide()
add_title(s, "Interpretabilidad: ¿de dónde viene la señal?")
add_image(s, os.path.join(EJ2, "reputation_tag_check.png"), 0.3, 1.3, height=4.2)
add_table(s, ["variante", "PR-AUC", "ROC-AUC"],
          [["original", "0,715", "0,961"], ["sin tag de title", "0,684", "0,944"],
           ["sin tag + sin frase de description", "0,142", "0,533"]],
          7.6, 1.5, 5.2, 1.6, size=12)
add_params_box(s, [
    "Modelo fusión semilla 0, sin reentrenar",
    "evaluado sobre 3 versiones del texto de valid",
], 7.6, 3.3, 5.2, 1.1)
add_conclusion_box(s,
    "description repite la misma señal de reputación con otra frase (95,4% de las "
    "filas) — no son 2 señales independientes. Confirmado con causalidad: recién al "
    "sacar ambas el PR-AUC se derrumba (por debajo del baseline tabular).",
    0.3, 5.0, 12.5, 2.0)

# ============================================================
# Estudio de ablacion
# ============================================================
s = slide()
add_title(s, "Estudio de ablación: arquitectura y features dudosas", size=26)
add_image(s, os.path.join(EJ2, "ablation.png"), 0.3, 1.3, height=4.1)
add_params_box(s, [
    "1 dial por vez desde la config base (64/4 heads/2 layers/128)",
    "1 sola semilla por variante (exploratorio)",
], 7.4, 1.4, 5.4, 1.1)
add_table(s, ["variante", "PR-AUC", "brecha train−valid"],
          [["base (64/4/2/128)", "0,716", "+0,192"], ["d_model=32", "0,779", "+0,128"],
           ["d_model=96", "0,710", "+0,239"], ["layers=1", "0,762", "+0,172"],
           ["sin country_of_origin", "0,763", "+0,183"], ["sin nutrition_score", "0,740", "+0,210"]],
          7.4, 2.6, 5.4, 2.9, size=11)
add_conclusion_box(s,
    "d_model=32 (menos de la mitad de parámetros) dio el mejor PR-AUC y menor brecha. "
    "d_model=96 fue el peor, con overfitting temprano confirmado. heads/layers dentro "
    "del ruido entre semillas (±0,017). country_of_origin/nutrition_score: sacarlas no "
    "empeora — confirma el EDA univariado.",
    0.3, 5.6, 12.5, 1.7)

# ============================================================
# Configuracion final
# ============================================================
s = slide()
add_title(s, "Configuración final: d_model=32 con 3 semillas")
add_table(s, ["semilla", "PR-AUC (valid)", "ROC-AUC (valid)", "brecha"],
          [["0", "0,779", "0,965", "+0,128"], ["1", "0,776", "0,970", "+0,117"],
           ["2", "0,743", "0,970", "+0,136"], ["media ± std", "0,766 ± 0,020", "0,968 ± 0,003", "+0,127"]],
          0.7, 1.6, 11.9, 2.2, size=15)
add_bullets(s, [
    "Comparado contra la base (d_model=64: 0,732 ± 0,017): mejor en media (+0,034 "
    "PR-AUC) con menos de la mitad de parámetros (45.569 vs. 102.337).",
    "Menor brecha train-valid (+0,127 vs. +0,156): sobreajusta menos.",
    "Las distribuciones se solapan parcialmente — la ventaja es consistente pero no "
    "abrumadora.",
], 0.7, 4.2, 11.9, 2.2, size=15)
add_conclusion_box(s,
    "achicar el modelo no empeora, y probablemente mejora. Configuración final: "
    "d_model=32, 4 heads, 2 encoders, MLP interno 128, fusión tardía.",
    0.7, 6.4, 11.9, 0.9)

# ============================================================
# BTR por busqueda
# ============================================================
s = slide()
add_title(s, "Del bought por fila al BTR por búsqueda")
add_image(s, os.path.join(EJ2, "btr_valid.png"), 0.4, 1.4, height=4.3)
add_table(s, ["métrica", "valor"],
          [["BTR real medio", "0,132"], ["BTR predicho medio", "0,135"],
           ["MAE", "0,067"], ["sesgo", "+0,003"],
           ["correlación (Pearson)", "0,764"], ["correlación de rangos (Spearman)", "0,742"]],
          7.5, 1.5, 5.3, 2.6, size=12)
add_params_box(s, [
    "Config final (d_model=32, semilla 0), 302 queries de valid",
    "BTR = promedio de bought / de predicciones por query_id",
], 7.5, 4.2, 5.3, 1.2)
add_conclusion_box(s,
    "bien calibrado en promedio (sesgo +0,003). Ordena bien las búsquedas (Spearman "
    "0,742) — importa más ordenar que acertar el valor exacto, coherente con no usar "
    "threshold.",
    0.4, 6.1, 12.5, 1.2)

# ============================================================
# Evaluacion final en test
# ============================================================
s = slide()
add_title(s, "Evaluación final en test")
add_table(s, ["semilla", "split", "PR-AUC", "ROC-AUC", "MAE del BTR"],
          [["0", "valid", "0,779", "0,965", "0,067"], ["0", "test", "0,782", "0,959", "0,067"],
           ["1", "valid", "0,749", "0,968", "0,071"], ["1", "test", "0,771", "0,962", "0,072"],
           ["2", "valid", "0,742", "0,970", "0,074"], ["2", "test", "0,733", "0,962", "0,072"],
           ["media", "valid", "0,756 ± 0,020", "0,968 ± 0,003", "0,071 ± 0,004"],
           ["media", "test", "0,762 ± 0,026", "0,961 ± 0,002", "0,070 ± 0,003"]],
          0.7, 1.4, 11.9, 4.3, size=12)
add_params_box(s, [
    "Test evaluado 1 sola vez, config congelada, época 20 fija (no la mejor de valid)",
], 0.7, 5.9, 11.9, 0.5)
add_conclusion_box(s,
    "sin degradación de valid a test (incluso levemente superior) — valida el "
    "procedimiento de 3 particiones y la elección de 70/15/15, no solo el modelo.",
    0.7, 6.5, 11.9, 0.8)

# ============================================================
# Desafios encontrados
# ============================================================
s = slide()
add_title(s, "Desafíos encontrados")
add_bullets(s, [
    "Agregar una medición cambió lo medido: iterar un DataLoader consume el generador "
    "de números aleatorios global de PyTorch, aunque sea con shuffle=False. Evaluar "
    "sobre train entre épocas corría el RNG y desalineaba dropout/shuffle posteriores. "
    "Solución: guardar y restaurar torch.get_rng_state() alrededor de esa evaluación.",
    "Los resultados del bloque Transformer dependen de la versión de PyTorch: el "
    "baseline tabular (solo capas densas) reproduce exacto entre entornos, pero "
    "nn.TransformerEncoder cambió de implementación interna entre versiones — la "
    "semilla sola no garantiza reproducibilidad, también hay que fijar versiones de "
    "librerías.",
], 0.7, 1.6, 11.9, 5.0, size=17)

# ============================================================
# Ejercicio 3
# ============================================================
s = slide()
add_title(s, "Ejercicio 3: Personalización")
add_bullets(s, [
    "Ejercicio teórico: cómo modificar la solución para incluir personalización de "
    "usuario al definir el BTR.",
    "Pendiente de desarrollar — no se trabajó todavía en este punto.",
], 0.7, 2.5, 11.9, 2.0, size=20)

# ============================================================
# Conclusiones generales
# ============================================================
s = slide()
add_title(s, "Conclusiones generales")
add_bullets(s, [
    "El Transformer aporta muchísimo sobre el baseline tabular (+0,55 PR-AUC) — no es un "
    "requisito formal sin sustancia, hay señal genuina en el texto.",
    "Esa señal viene, casi en su totalidad, de un tag de reputación redundante entre "
    "title y description — identificado en el EDA y confirmado causalmente en el "
    "Ejercicio 2.",
    "d_model chico (32) rindió mejor que uno más grande (64/96), con menos overfitting "
    "y menos de la mitad de parámetros — \"empezar chico\" (sugerido en la consigna) "
    "resultó ser lo mejor, no solo lo más barato.",
    "El BTR agregado por búsqueda queda bien calibrado y ordena bien, sin necesidad de "
    "threshold.",
    "Resultado final en test (nunca usado en ninguna decisión): PR-AUC 0,762 ± 0,026, "
    "ROC-AUC 0,961 ± 0,002, sin degradación respecto de valid.",
], 0.7, 1.5, 11.9, 5.3, size=16)

# ============================================================
# Fin
# ============================================================
s = slide()
add_text(s, "Gracias", 0.5, 3.3, 12.3, 1.2, size=44, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

out_path = os.path.join(BASE, "presentacion_tp1.pptx")
prs.save(out_path)
print("Guardado en", out_path, "-", len(prs.slides), "slides")
